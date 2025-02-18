from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Define a function to add slides with titles and content
def add_slide(prs, title, content, bullet_points=False):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    slide.shapes.title.text = title
    if bullet_points:
        content_box = slide.placeholders[1]
        for point in content:
            p = content_box.text_frame.add_paragraph()
            p.text = point
            p.level = 0
    else:
        slide.placeholders[1].text = content

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Virtual Machines: The Future of Flexible Computing"
slide.placeholders[1].text = "Run anything, anywhere."

# Slide 2: What is a Virtual Machine?
add_slide(
    prs,
    "What is a Virtual Machine?",
    "A virtual machine is a software emulation of a computer system.\n"
    "It runs like a physical machine but is virtualized inside another system.",
)

# Slide 3: Why Use Virtual Machines?
add_slide(
    prs,
    "Why Use Virtual Machines?",
    [
        "Resource Efficiency: Use one server for multiple systems.",
        "Flexibility: Run different OSes on the same hardware.",
        "Security: Isolated environments prevent cross-contamination.",
        "Portability: Easily migrate VMs across devices.",
    ],
    bullet_points=True,
)

# Slide 4: Types of Virtual Machines
add_slide(
    prs,
    "Types of Virtual Machines",
    [
        "System VMs: Virtualize an entire OS.",
        "Process VMs: Run a single application (e.g., Java Virtual Machine).",
    ],
    bullet_points=True,
)

# Slide 5: How Do VMs Work?
add_slide(
    prs,
    "How Do VMs Work?",
    "Key Components:\n"
    "- Hypervisor (Type 1 & Type 2)\n"
    "- Virtualized hardware\n\n"
    "Process: Hardware → Hypervisor → VM → OS/Application",
)

# Slide 6: Benefits of Virtual Machines
add_slide(
    prs,
    "Benefits of Virtual Machines",
    [
        "Cost Savings: Maximize hardware utilization.",
        "Disaster Recovery: Snapshots for quick restoration.",
        "Development & Testing: Isolated and reproducible environments.",
    ],
    bullet_points=True,
)

# Slide 7: Challenges of Virtual Machines
add_slide(
    prs,
    "Challenges of Virtual Machines",
    [
        "Performance Overhead.",
        "Complexity in Management.",
        "Limited Direct Access to Hardware.",
    ],
    bullet_points=True,
)

# Slide 8: Virtual Machines vs Containers
add_slide(
    prs,
    "Virtual Machines vs Containers",
    "Virtual Machines:\n"
    "- Full OS, heavy, slower to start.\n\n"
    "Containers:\n"
    "- Shared OS, lightweight, faster.",
)

# Slide 9: Use Cases for Virtual Machines
add_slide(
    prs,
    "Use Cases for Virtual Machines",
    [
        "Running Legacy Software.",
        "Development & Testing Environments.",
        "Cloud Computing Infrastructure.",
        "Secure Browsing or Sandbox Environments.",
    ],
    bullet_points=True,
)

# Slide 10: Future of Virtual Machines
add_slide(
    prs,
    "Future of Virtual Machines",
    [
        "Integration with Cloud Services.",
        "Enhanced Performance via Hardware Virtualization.",
        "AI-powered VM Optimization.",
    ],
    bullet_points=True,
)

# Slide 11: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
slide.placeholders[1].text = "Virtual Machines: Empowering the Digital World."

# Save the presentation
prs.save("Virtual_Machines_Presentation.pptx")

print("Presentation saved as 'Virtual_Machines_Presentation.pptx'")

